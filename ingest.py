"""Build TF-IDF index over the PIF training corpus with adaptive chunking.

corpus_files/
  methodology/     official PDFs — paragraph chunks 800 chars + neighbor expansion at retrieval
  training/        training decks — paragraph chunks 1200 chars, slide as unit
  case_studies/    case study decks — 600 char chunks, boilerplate-stripped, keyword-tagged
  district_data/   per-district-sector .txt files — one file = one chunk (already small)
  vision_documents/ state/district/constituency/mandal vision plans — paragraph chunks 800 chars
"""
import os
import pickle
import re

import docx
import pdfplumber
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer

CORPUS_DIR = os.path.join(os.path.dirname(__file__), "corpus_files")
INDEX_PATH = os.path.join(os.path.dirname(__file__), "index.pkl")

CHUNK_SIZES = {
    "methodology":   800,
    "training":     1200,
    "case_studies":  600,
    "district_data": 99999,  # whole file = one chunk
    "vision_documents": 2000,  # larger chunks keep the embedding count (and build time) manageable
}
OVERLAP = 150

CASE_STUDY_TAGS = {
    "Srikakulam_Blue_Economy.txt": "fisheries aquaculture shrimp fish marine coastal blue economy",
    "Nellore_Shrimp_Processing.txt": "fisheries aquaculture shrimp seafood processing export",
    "Biofloc_Tilapia_CaseStudy.txt": "fisheries aquaculture fish tilapia biofloc pond farming",
    "Paddy_Fish_Integrated_Farming_Case_Study_AP.txt": "fisheries aquaculture fish paddy integrated farming",
    "East_Godavari_Coconut_Coir.txt": "coconut coir horticulture agro processing waste to wealth",
    "Nellore_Ethanol_Potential.txt": "ethanol biofuel sugarcane maize distillery energy manufacturing",
    "Banana_Processing_Case_Study.txt": "banana horticulture agro processing waste to wealth pseudo-stem",
    "Kumarakom_Responsible_Tourism_Case_Study.txt": "tourism hospitality services responsible tourism ecotourism",
    "Shenzhen_Port_Led_Manufacturing.txt": "manufacturing industry port logistics special economic zone china",
    "Sahyadri_Replication_Playbook.txt": "grapes horticulture fpo fpc farmer producer company cold chain",
    "Sahyadri_V2.txt": "chilli spices agriculture fpo fpc value chain export prakasam",
    "Chetna_FPO_Lessons.txt": "organic cotton textiles fpo shg cooperative farmer producer company",
    "Morbi_Ceramics_Industry.txt": "ceramics tiles manufacturing industry cluster gujarat",
    "Tiruppur_Case_Study_Updated.txt": "textiles knitwear garments manufacturing export cluster zero liquid discharge",
}

BOILERPLATE_PATTERNS = [
    re.compile(r"^CASE STUDY\s*\|.*$", re.IGNORECASE | re.MULTILINE),
    re.compile(r".*Training deck for District.*Mandal Level Officials.*$", re.IGNORECASE | re.MULTILINE),
    re.compile(r"^Editable PowerPoint\s*\|.*slides.*$", re.IGNORECASE | re.MULTILINE),
    re.compile(r"^Source:.*$", re.IGNORECASE | re.MULTILINE),
    re.compile(r".*\|\s*Case study for AP officials.*$", re.IGNORECASE | re.MULTILINE),
]


def strip_boilerplate(text):
    for pat in BOILERPLATE_PATTERNS:
        text = pat.sub("", text)
    return text


def extract_docx_units(path):
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return [{"page": None, "text": "\n".join(parts)}]


def extract_pptx_units(path):
    prs = Presentation(path)
    units = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t.strip():
                        slide_text.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    slide_text.append(" | ".join(c.text.strip() for c in row.cells))
        if slide_text:
            units.append({"page": i, "text": "\n".join(slide_text)})
    return units


def extract_pdf_units(path):
    units = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            t = page.extract_text() or ""
            if t.strip():
                units.append({"page": i, "text": t})
    return units


def extract_txt_units(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return [{"page": None, "text": f.read()}]


def extract_units(path):
    ext = path.lower().rsplit(".", 1)[-1]
    try:
        if ext == "docx":
            return extract_docx_units(path)
        if ext == "pptx":
            return extract_pptx_units(path)
        if ext == "pdf":
            return extract_pdf_units(path)
        if ext == "txt":
            return extract_txt_units(path)
    except Exception as e:
        print(f"  ! failed to extract {path}: {e}")
    return []


def sub_split(text, page, source, max_chars):
    """Split at paragraph boundaries, keeping the same page label."""
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [{"source": source, "page": page, "text": text}]
    paragraphs = text.split("\n\n")
    chunks, current = [], ""
    for para in paragraphs:
        if len(current) + len(para) > max_chars and current:
            chunks.append({"source": source, "page": page, "text": current.strip()})
            current = current[-OVERLAP:] + "\n\n" + para
        else:
            current = (current + "\n\n" + para).strip()
    if current.strip():
        chunks.append({"source": source, "page": page, "text": current.strip()})
    return chunks


def process_folder(folder_name, folder_path, all_chunks):
    max_chars = CHUNK_SIZES.get(folder_name, 1200)
    is_case_studies = folder_name == "case_studies"
    is_district = folder_name == "district_data"

    for root, _, files in os.walk(folder_path):
        for fname in sorted(files):
            if fname.startswith("."):
                continue
            ext = fname.lower().rsplit(".", 1)[-1]
            if ext not in ("docx", "pptx", "pdf", "txt"):
                continue
            path = os.path.join(root, fname)
            rel = os.path.relpath(path, CORPUS_DIR)

            units = extract_units(path)
            file_chunks = []
            tags = CASE_STUDY_TAGS.get(fname) if is_case_studies else None

            for u in units:
                text = strip_boilerplate(u["text"]) if is_case_studies else u["text"]
                for chunk in sub_split(text, u["page"], rel, max_chars):
                    if tags:
                        tag_block = " ".join([tags] * 4)
                        chunk["text"] = f"Keywords: {tag_block}\n\n{chunk['text']}"
                    # district chunks get their folder as metadata for filtering
                    chunk["folder"] = folder_name
                    file_chunks.append(chunk)

            print(f"  {rel}: {len(units)} units → {len(file_chunks)} chunks")
            all_chunks.extend(file_chunks)


FOLDERS = ["methodology", "training", "case_studies", "district_data", "vision_documents"]


def build_all_chunks():
    """Extract + chunk every corpus folder. Shared by TF-IDF and embedding indexers."""
    all_chunks = []
    for folder in FOLDERS:
        path = os.path.join(CORPUS_DIR, folder)
        if not os.path.isdir(path):
            print(f"Skipping {folder} (not found)")
            continue
        print(f"\n[{folder.upper()}] max_chars={CHUNK_SIZES[folder]}")
        process_folder(folder, path, all_chunks)
    return all_chunks


def main():
    all_chunks = build_all_chunks()

    if not all_chunks:
        print("No chunks extracted.")
        return

    texts = [c["text"] for c in all_chunks]
    vectorizer = TfidfVectorizer(stop_words="english", max_features=20000, ngram_range=(1, 2))
    matrix = vectorizer.fit_transform(texts)

    with open(INDEX_PATH, "wb") as f:
        pickle.dump({"chunks": all_chunks, "vectorizer": vectorizer, "matrix": matrix}, f)

    by_folder = {}
    for c in all_chunks:
        by_folder[c["folder"]] = by_folder.get(c["folder"], 0) + 1
    print(f"\nTotal: {len(all_chunks)} chunks")
    for k, v in by_folder.items():
        print(f"  {k}: {v}")
    print(f"Index saved → {INDEX_PATH}")


if __name__ == "__main__":
    main()
